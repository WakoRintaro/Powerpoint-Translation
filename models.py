import csv
from pptx import Presentation
import unicodedata
import requests
import json
import os

class Transpptx:
    """
    Class for Powerpoint Translation
    
    Parameters
    ----------
    pptx: string
        File Path to target pptx
    saveto: string
        File Path to translated pptx
    csv: string
        File Path to the dictionary csv
    """
    wordbank = {}
    num_query_api = 0
    num_query_bank = 0
    
    def __init__(self, pptx, saveto, csv):
        """
        Create an instance
        
        ex)
            savetoに指定したパスとpptxのファイル名の組み合わせが保存先
            pptx = /home/xxx/filename.pptx
            saveto = /home/xxx/translated/
            self.path_translated_pptx = /home/xxx/translated/filename.pptx
        """
        self.path_target_pptx = pptx
        self.path_translated_pptx = saveto + os.path.basename(self.path_target_pptx)
        self.path_dictionary = csv
        print(self.path_translated_pptx)

    def __load_dictionary(self):
        """
        Load Dictionary CSV
        """
        self.dict_fieldnames = ['jpn','rubi','en','th']
        with open(self.path_dictionary, "r", newline='', encoding='utf-8') as f:
            reader = csv.DictReader(f, fieldnames=self.dict_fieldnames)
            for row in reader:
                self.wordbank[row['jpn']] = {"rubi":row['rubi'], "en":row['en'], "th":row['th']}
    
    def __save_dictionary(self):
        """
        Save Dictionary CSV
        """
        w = csv.writer(open(self.path_dictionary, "w", newline='',encoding='utf-8'))
        for key, val in self.wordbank.items():
            w.writerow([key, val['rubi'], val['en'], val['th']])

    def pptx_translation(self):
        """
        Execute Translation of pptx and Save
        """
        # Load Dictionary CSV
        self.__load_dictionary()

        prs = Presentation(self.path_target_pptx)
        prs.notes_master

        for ns, slide in enumerate(prs.slides):
            text_slide = ""
            # print("Slide No." + str(ns+1))
            for nsh, shape in enumerate(slide.shapes):
                text_shape = ""
                if not shape.has_text_frame:
                    continue
                for np, paragraph in enumerate(shape.text_frame.paragraphs):
                    text_parag = ""
                    for rs, run in enumerate(paragraph.runs):
                        text_parag += run.text
                    text_shape += self.__text_translate(text_parag)
                    # text_shape += text_parag
                text_slide += text_shape
            
            # Note Slide
            note_slide = slide.notes_slide
            note_slide.notes_placeholder
            # print("NOTE -> ",slide.has_notes_slide)
            # print(note_slide.placeholders)

            # TextFrame in Note Slide
            text_frame = note_slide.notes_text_frame
            # print(text_frame)

            # Text in TextFrame
            # When it doesn't have text
            try:
                text_frame.text = text_slide
            except AttributeError as e:
                try:
                    text_frame.text = text_slide
                    print("Success (after the 2nd trial")
                except AttributeError as e2:
                    print("Failed (after the 2nd trial)")
                    print(e2)
                print(e)

        # Save Translated pptx
        print("save", self.path_translated_pptx)
        prs.save(self.path_translated_pptx)

        # Save Dictionray CSV
        self.__save_dictionary()

    def __text_translate(self, text):
        """
        Translate Japanese to Rubi, English and Thai
        Return empty text if translation is not needed
        """

        # When the text has no characters
        if len(text) == 0:
            return ""
        
        # When the text doesn't have Kanji
        if not self.__has_japanese(text):
            return ""

        # Check if the text exsits in WORD_BANK
        if self.__is_in_bank(text):
            trans = self.__translate_by_bank(text)
        else:
            trans = self.__translate_by_api(text)

        # When the api fails to translate
        if trans is None:
            return ""   
        return self.__format_memo(trans)

    def __format_memo(self, trans):
        """
        format translation
        ex:
            南京錠 / (なんきんじょう)
            padlock
            กุญแจ
        """
        rtn_text = trans['jpn'] + " / " + trans['rubi'] + "\n"
        rtn_text += trans['en'] + "\n" + trans['th'] + "\n\n"
        return rtn_text

    def __has_japanese(self, text):
        """
        Return TRUE if the text has at least one Kanji character
        """
        for ch in text:
            name = unicodedata.name(ch, 'Undefined')
            if "CJK UNIFIED" in name:
                return True
        return False

    def __is_in_bank(self, text):
        """
        Return if the wordbank has cashe
        """
        return text in self.wordbank
    
    def __translate_by_bank(self, text):
        """
        Return translation by WORD BANK
        """
        self.num_query_bank += 1
        
        trans = {}
        trans['jpn'] = text
        trans.update(self.wordbank[text])
        return trans

    def __add_to_bank(self, trans):
        """
        Add translation to WORDBANK
        """
        # global WORD_BANK
        self.wordbank[trans['jpn']] = {"rubi":trans['rubi'], "en":trans['en'], "th":trans['th']}

    def __translate_by_api(self, text):
        """
        Return translation by API
        """
        self.num_query_api += 1

        # Request translation to Google Apps
        url = "https://script.google.com/macros/s/AKfycbzFsUuMuFmZ7mCfWMmGYrqZJp_XIDHuUz9JcyZx0-oOIZHXEQ/exec"
        rr = requests.get(url, params={'text':text}, timeout=(3.0, 7.5) )
        rr.encoding = rr.apparent_encoding
        # print(rr.status_code) 
        # When the response is against the JSON formatting
        try:
            trans = json.loads(rr.text)
        except json.JSONDecodeError as e:
            print("An Error Occured in API")
            # print(e)
            # print("REQUEST:", text)
            # print("RECIEVED", rr.text)
            return None
        self.__add_to_bank(trans)
        return trans
